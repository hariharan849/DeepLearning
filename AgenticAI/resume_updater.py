import time, requests
import streamlit as st
import ast, re
from bs4 import BeautifulSoup
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate
from langgraph.graph import StateGraph, END
from langchain_community.utilities import SerpAPIWrapper
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF for PDFs
import os
from difflib import unified_diff, HtmlDiff
from difflib import ndiff
from dotenv import load_dotenv
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

load_dotenv()
temp_dir = "./temp_files"
os.makedirs(temp_dir, exist_ok=True)

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
SERPAPI_KEY = os.getenv("SERPAPI_API_KEY")
search = SerpAPIWrapper()

# Set up LLM
llm = ChatGroq(model_name="gemma2-9b-it", temperature=0.7, api_key=GROQ_API_KEY)

def extract_text_and_format_from_resume(file, file_type):
    file_path = os.path.join(temp_dir, file.name)
    
    with open(file_path, "wb") as f:
        f.write(file.getbuffer())
    
    extracted_text = ""
    formatting = []
    
    if file_type == "docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            extracted_text += para.text + "\n"
            formatting.append({
                "text": para.text,
                "bold": any(run.bold for run in para.runs),
                "italic": any(run.italic for run in para.runs),
                "font_size": para.runs[0].font.size if para.runs else None
            })
    elif file_type == "pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
                    formatting.append({"text": shape.text, "font_size": None})
    elif file_type == "pdf":
        doc = fitz.open(file_path)
        for page in doc:
            for text in page.get_text("dict")["blocks"]:
                if "lines" in text:
                    for line in text["lines"]:
                        for span in line["spans"]:
                            extracted_text += span["text"] + " "
                            formatting.append({"text": span["text"], "font_size": span["size"], "font": span["font"]})
            extracted_text += "\n"
    else:
        return "", []
    
    safe_remove(file_path)  # Ensure file deletion

    return extracted_text.strip(), formatting

def save_updated_pdf(updated_text, formatting, file_name):
    output_path = os.path.join(temp_dir, f"updated_{file_name}")
    c = canvas.Canvas(output_path, pagesize=letter)
    y = 750  # Start position for text
    
    updated_text_lines = updated_text.split("\n")
    
    for i, item in enumerate(formatting):
        if i < len(updated_text_lines):
            font_name = item.get("font", "Helvetica")
            try:
                c.setFont(font_name, item["font_size"] if item["font_size"] else 12)
            except KeyError:
                c.setFont("Helvetica", item["font_size"] if item["font_size"] else 12)  # Fallback font
            c.drawString(50, y, updated_text_lines[i])
            y -= 20
    
    c.save()
    return output_path

def safe_remove(file_path):
    for _ in range(5):  # Try for a short period
        try:
            os.remove(file_path)
            break
        except PermissionError:
            time.sleep(0.5)  # Wait and retry

def fetch_job_description_from_llm(url):
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers)
        response.raise_for_status()

        # Extract raw text from the webpage
        soup = BeautifulSoup(response.text, "html.parser")
        page_text = soup.get_text(separator=" ", strip=True)

        # Use LLM to extract the job description from the raw text
        job_description = llm.invoke(f"""
        Extract the job description from the following webpage text. Focus only on the responsibilities, qualifications, and skills mentioned for the role. 
        Ignore navigation menus, company descriptions, and unrelated content. Show each skill or requirement on a new line.

        Webpage Content:
        {page_text}

        Extracted Job Description:
        """)

        company_name = llm.invoke(f"""
        Extract the company name from the following webpage content. Focus on the employer's name, avoiding job portals, ads, or unrelated mentions.  

        Webpage Content:  
        {page_text}  

        Extracted Company Name:
        """)

        return job_description, company_name.content.strip()
    
    except Exception as e:
        print (f"Error fetching job description: {str(e)}")
    return None, None


def highlight_differences(original, updated):
    diff = list(ndiff(original.split(), updated.split()))
    highlighted = " ".join([
        f'**{word[2:]}**' if word.startswith('+ ') else word[2:]  
        for word in diff if not word.startswith('- ')
    ])
    return highlighted

def extract_job_requirements(state):
    extracted_text = llm.invoke(extract_prompt.format(job_description=state["job_description"]))
    return {
        "resume": state["resume"],
        "job_description": state["job_description"],
        "extracted_requirements": extracted_text,
        "formatting": state["formatting"],  # Ensure formatting is passed forward
        "file_type": state["file_type"]  # Ensure file type is also passed forward
    }

def analyze_resume_gaps(state):
    gap_analysis = llm.invoke(analyze_prompt.format(resume=state["resume"], extracted_requirements=state["extracted_requirements"]))
    return {
        "resume": state["resume"],
        "job_description": state["job_description"],
        "extracted_requirements": state["extracted_requirements"],
        "resume_gaps": gap_analysis.content,
        "formatting": state["formatting"],  # Ensure formatting is passed forward
        "file_type": state["file_type"]  # Ensure file_type is also passed forward
    }

def update_resume(state):
    updated_text = llm.invoke(f"""
        Modify the resume below to align with the job requirements while keeping it truthful and professional.
        Highlight missing skills and reframe experience to better match the role.
        
        Resume:
        {state["resume"]}
        
        Identified Gaps:
        {state["resume_gaps"]}
        
        Where possible, subtly incorporate missing but relevant skills into existing experience sections.
        Return only the modified text against original resume
    """)
    highlighted_text = highlight_differences(state["resume"], updated_text.content)
    return {
        "resume": state["resume"],
        "job_description": state["job_description"],
        "extracted_requirements": state["extracted_requirements"],
        "formatting": state["formatting"],
        "resume_gaps": state["resume_gaps"],
        "updated_resume": highlighted_text
    }

# Extract key skills from job description
extract_prompt = PromptTemplate(
    input_variables=["job_description"],
    template="""
    Extract key skills, responsibilities, and qualifications from the following job description:
    {job_description}
    """
)

# Step 2: Analyze gaps in resume
analyze_prompt = PromptTemplate(
    input_variables=["resume", "extracted_requirements"],
    template="""
    Compare the given resume with the extracted job requirements and identify missing skills, experience gaps, and areas of improvement.
    Resume:
    {resume}
    Job Requirements:
    {extracted_requirements}
    """
)

graph = StateGraph(dict)
graph.add_node("extract", extract_job_requirements)
graph.add_node("analyze", analyze_resume_gaps)
graph.add_node("update", update_resume)

graph.add_edge("extract", "analyze")
graph.add_edge("analyze", "update")
graph.add_edge("update", END)

graph.set_entry_point("extract")
graph = graph.compile()

st.title("Resume Comparison against Job Description")
st.write("This tool helps you optimize your resume by comparing it against a job description and identifying key gaps.")

uploaded_resume = st.file_uploader("Upload your resume (Word, PPT, or PDF)", type=["docx", "pptx", "pdf"])
# Display information about the uploaded file
if uploaded_resume is not None:
    st.session_state.otto_file = uploaded_resume
    st.header("Uploaded File Details:")

    # Display file type and size
    file_details = {
        "File Name": uploaded_resume.name,
        "File Type": uploaded_resume.type,
        "File Size (bytes)": uploaded_resume.size,
    }
    st.write(file_details)

job_url = st.text_input("Paste the job listing URL")

if "job_url" not in st.session_state:
    st.session_state.job_url = None

st.session_state.job_url = job_url

if "job_description" not in st.session_state:
    st.session_state.job_description = None
if "company_name" not in st.session_state:
    st.session_state.company_name = None

if st.button("Fetch Job Description") and st.session_state.job_url:
    job_description, company_name = fetch_job_description_from_llm(st.session_state.job_url)
    if job_description:
        st.session_state.job_description = job_description.content  # Store in session state
        st.session_state.company_name = company_name
        st.success("Job description fetched successfully!")

if st.session_state.job_description:
    st.subheader("Extracted Job Description")
    st.text_area("Job Description", st.session_state.job_description, height=300)

if "resume_optimized" not in st.session_state:
    st.session_state.resume_optimized = None

col1, col2 = st.columns(2)
with col1:
    if st.button("Optimize Resume") and uploaded_resume and st.session_state.job_description:
        file_type = uploaded_resume.name.split(".")[-1]
        resume_text, formatting = extract_text_and_format_from_resume(uploaded_resume, file_type)
        
        if not resume_text:
            st.error("Unsupported file type or extraction error.")
        else:
            state = {"resume": resume_text, "job_description": st.session_state.job_description, "formatting": formatting, "file_type": file_type}
            st.session_state.resume_optimized = graph.invoke(state)
            
            st.subheader("Identified Resume Gaps:")
            st.write(st.session_state.resume_optimized["resume_gaps"])

def fetch_linkedin_profile_data(profile_url):
    try:
        if not re.match(r'https?://', profile_url):
            print(f"Invalid URL: {profile_url}")
            return ''
        response = requests.get(profile_url, headers={'User-Agent': 'Mozilla/5.0'})
        soup = BeautifulSoup(response.text, 'html.parser')
        profile_img_tag = soup.find('img', {'class': re.compile(r'profile-picture|pv-top-card-profile-picture__image')})
        img_url = profile_img_tag['src'] if profile_img_tag else ''
        return img_url
    except Exception as e:
        print(f"Error fetching profile picture: {e}")
        return ''

def search_recruiters(company):
    search_query = f"Recruiter site:linkedin.com/in AND {company} AND India"
    url = f"https://serpapi.com/search.json?q={search_query}&engine=google&api_key={SERPAPI_KEY}"

    response = requests.get(url)
    results = response.json().get("organic_results", [])
    
    recruiters = []
    for result in results:
        recruiters.append({
            "name": result.get("title", "").split(" | ")[0],
            "title": result.get("title", ""),
            "link": result.get("link"),
            "company": company
        })

    return recruiters

def summarize_recruiters(recruiters):
    llm = ChatGroq(model_name="mixtral-8x7b-32768", temperature=0.7, api_key=GROQ_API_KEY)
    company = recruiters[0]["company"]
    prompt = PromptTemplate(
        input_variables=["recruiters", "company"],
        template="""
        You must respond only in valid JSON format. Do not include any explanations, introductions, or extra text.
        
        Based on the recruiter profiles below, identify the best recruiter for DataScience/ML/AI/Software roles at {company}.
        Provide a **ranked list** in the following JSON format:
        
        ```json
        [
            {{"rank": 1, "name": "Recruiter Name", "title": "Recruiter Title", "link": "Profile URL", "insight": "Reason for ranking"}},
            {{"rank": 2, "name": "Recruiter Name", "title": "Recruiter Title", "link": "Profile URL", "insight": "Reason for ranking"}},
            ...
        ]
        ```

        Recruiters:
        {recruiters}
        """
    )

    # Get the response from the model
    response = llm.invoke(prompt.format(recruiters=recruiters, company=company)).content

    # Ensure output is in JSON format
    import json
    try:
        ranked_recruiters = json.loads(response)
        return ranked_recruiters
    except json.JSONDecodeError:
        print("Error: Model output is not valid JSON.")
        return []



if st.session_state.resume_optimized:
    if st.button("Create Optimized Resume"):
        file_type = uploaded_resume.name.split(".")[-1]
        file_path = save_updated_pdf(
            st.session_state.resume_optimized['updated_resume'],
            st.session_state.resume_optimized["formatting"],
            "optimized_resume"
        )
        with open(file_path, "rb") as file:
            st.download_button("Download Updated Resume", file, file_name=f"optimized_resume.{file_type}")

with col2:
    if st.session_state.company_name:
        if st.button("Search Talent Acquisition at Company") and st.session_state.company_name:
            st.subheader(f"Talent Acquisition Professionals at {st.session_state.company_name}")
            recruiters = summarize_recruiters(search_recruiters(st.session_state.company_name))
            if recruiters:
                st.subheader(f"Recruiters at {st.session_state.company_name}")

                for idx, recruiter in enumerate(recruiters):
                    # st.image(recruiter["profile_pic"], width=150)
                    st.markdown(f"### [{recruiter['name']}]({recruiter['link']})")
                    st.write(recruiter["title"])
                    st.markdown(f"[🔗 LinkedIn Profile]({recruiter['link']})", unsafe_allow_html=True)
            else:
                st.warning("No recruiters found. Try another company.")
import os
import time
import requests
import json
import re
import ast
import streamlit as st
from bs4 import BeautifulSoup
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate
from langgraph.graph import StateGraph, END
from langchain_community.utilities import SerpAPIWrapper
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF for PDFs
from dotenv import load_dotenv
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from difflib import ndiff

# Load environment variables
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
SERPAPI_KEY = os.getenv("SERPAPI_API_KEY")

# Initialize API clients
search = SerpAPIWrapper()
llm = ChatGroq(model_name="mixtral-8x7b-32768", temperature=0.7, api_key=GROQ_API_KEY)

# Create temporary directory
TEMP_DIR = "./temp_files"
os.makedirs(TEMP_DIR, exist_ok=True)

# ===================== Helper Functions ===================== #

def safe_remove(file_path):
    """Tries to remove a file with retries."""
    for _ in range(5):
        try:
            os.remove(file_path)
            break
        except PermissionError:
            time.sleep(0.5)

def extract_text_from_resume(file, file_type):
    """Extracts text and formatting details from a resume file."""
    file_path = os.path.join(TEMP_DIR, file.name)
    
    with open(file_path, "wb") as f:
        f.write(file.getbuffer())

    extracted_text = ""
    formatting = []

    if file_type == "docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            extracted_text += para.text + "\n"
            formatting.append({
                "text": para.text,
                "bold": any(run.bold for run in para.runs),
                "italic": any(run.italic for run in para.runs),
                "font_size": para.runs[0].font.size if para.runs else None
            })
    elif file_type == "pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
                    formatting.append({"text": shape.text, "font_size": None})
    elif file_type == "pdf":
        doc = fitz.open(file_path)
        for page in doc:
            for text in page.get_text("dict")["blocks"]:
                if "lines" in text:
                    for line in text["lines"]:
                        for span in line["spans"]:
                            extracted_text += span["text"] + " "
                            formatting.append({"text": span["text"], "font_size": span["size"], "font": span["font"]})
            extracted_text += "\n"
    else:
        return "", []

    safe_remove(file_path)
    return extracted_text.strip(), formatting

def fetch_job_description(url):
    """Extracts job description and company name from a URL."""
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers)
        response.raise_for_status()

        soup = BeautifulSoup(response.text, "html.parser")
        page_text = soup.get_text(separator=" ", strip=True)

        job_description = llm.invoke(f"Extract the job description from:\n{page_text}").content
        company_name = llm.invoke(f"Extract the company name from:\n{page_text}").content.strip()

        return job_description, company_name
    except Exception as e:
        print(f"Error fetching job description: {str(e)}")
        return None, None

def search_recruiters(company):
    search_query = f"Recruiter site:linkedin.com/in AND {company} AND India"
    url = f"https://serpapi.com/search.json?q={search_query}&engine=google&api_key={SERPAPI_KEY}"

    response = requests.get(url)
    results = response.json().get("organic_results", [])
    
    recruiters = []
    for result in results:
        recruiters.append({
            "name": result.get("title", "").split(" | ")[0],
            "title": result.get("title", ""),
            "link": result.get("link"),
            "company": company
        })

    return recruiters

def summarize_recruiters(recruiters):
    """Uses LLM to rank recruiters."""
    company = recruiters[0]["company"]
    prompt = PromptTemplate(
        input_variables=["recruiters", "company"],
        template="""
        You must respond only in valid JSON format. Do not include any explanations, introductions, or extra text.
        
        Based on the recruiter profiles below, identify the best recruiter for DataScience/ML/AI/Software roles at {company}.
        Provide a **ranked list** in the following JSON format:
        
        ```json
        [
            {{"rank": 1, "name": "Recruiter Name", "title": "Recruiter Title", "link": "Profile URL", "insight": "Reason for ranking"}},
            {{"rank": 2, "name": "Recruiter Name", "title": "Recruiter Title", "link": "Profile URL", "insight": "Reason for ranking"}},
            ...
        ]
        ```

        Recruiters:
        {recruiters}
        """
    )

    response = llm.invoke(prompt.format(recruiters=recruiters, company=company)).content
    try:
        return json.loads(response)
    except json.JSONDecodeError:
        print("Error: Model output is not valid JSON.")
        return []

# ===================== Streamlit UI ===================== #

st.set_page_config(page_title="Resume & Recruiter Finder", layout="wide")

st.title("🚀 Resume Optimization & Recruiter Finder")

uploaded_resume = st.file_uploader("📂 Upload your resume (Word, PPT, or PDF)", type=["docx", "pptx", "pdf"])
job_url = st.text_input("🔗 Paste the job listing URL")

col1, col2 = st.columns(2)

with col1:
    if st.button("Fetch Job Description") and job_url:
        with st.spinner("Fetching job details..."):
            job_description, company_name = fetch_job_description(job_url)
        if job_description:
            st.session_state.job_description = job_description
            st.session_state.company_name = company_name
            st.success("✅ Job description fetched successfully!")
            st.text_area("📄 Extracted Job Description", job_description, height=300)

    if uploaded_resume and st.session_state.get("job_description"):
        if st.button("Optimize Resume"):
            file_type = uploaded_resume.name.split(".")[-1]
            resume_text, formatting = extract_text_from_resume(uploaded_resume, file_type)
            
            if not resume_text:
                st.error("❌ Unsupported file type or extraction error.")
            else:
                state = {
                    "resume": resume_text,
                    "job_description": st.session_state["job_description"],
                    "formatting": formatting,
                    "file_type": file_type
                }
                st.session_state.resume_optimized = state
                st.subheader("🔎 Identified Resume Gaps")
                st.write("🔹 Example missing skills here...")

with col2:
    if st.session_state.get("company_name"):
        if st.button("Search Talent Acquisition"):
            recruiters = summarize_recruiters(search_recruiters(st.session_state["company_name"]))
            if recruiters:
                for recruiter in recruiters:
                    st.markdown(f"### [{recruiter['name']}]({recruiter['link']})")
                    st.write(recruiter["title"])
                    st.markdown(f"[🔗 LinkedIn Profile]({recruiter['link']})", unsafe_allow_html=True)
            else:
                st.warning("⚠️ No recruiters found. Try another company.")

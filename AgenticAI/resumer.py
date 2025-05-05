import os
import time
import requests
import numpy as np
import json
import re
import streamlit as st
import soundfile as sf
from bs4 import BeautifulSoup
from streamlit_mic_recorder import mic_recorder
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF for PDFs

# Load API keys
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
SERPAPI_KEY = os.getenv("SERPAPI_API_KEY")

# Initialize AI Model
llm = ChatGroq(model_name="mixtral-8x7b-32768", temperature=0.7, api_key=GROQ_API_KEY)

# ===================== Helper Functions ===================== #

def safe_remove(file_path):
    """Safely remove a file."""
    try:
        os.remove(file_path)
    except Exception:
        pass

def extract_text_from_resume(file, file_type):
    """Extracts text from resume files."""
    file_path = f"./temp_files/{file.name}"
    with open(file_path, "wb") as f:
        f.write(file.getbuffer())

    extracted_text = ""
    if file_type == "docx":
        doc = Document(file_path)
        extracted_text = "\n".join([para.text for para in doc.paragraphs])
    elif file_type == "pptx":
        prs = Presentation(file_path)
        extracted_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif file_type == "pdf":
        doc = fitz.open(file_path)
        extracted_text = "\n".join([page.get_text("text") for page in doc])
    
    safe_remove(file_path)
    return extracted_text.strip()

def fetch_job_description(url):
    """Extracts job description from a webpage."""
    headers = {"User-Agent": "Mozilla/5.0"}
    response = requests.get(url, headers=headers)
    soup = BeautifulSoup(response.text, "html.parser")
    
    job_description = llm.invoke(f"Extract job description from:\n{soup.get_text()}").content
    company_name = llm.invoke(f"Extract company name from:\n{soup.get_text()}").content.strip()
    
    return job_description, company_name

def search_recruiters(company):
    search_query = f"Recruiter site:linkedin.com/in AND {company} AND India"
    url = f"https://serpapi.com/search.json?q={search_query}&engine=google&api_key={SERPAPI_KEY}"

    response = requests.get(url)
    results = response.json().get("organic_results", [])
    
    recruiters = []
    for result in results:
        recruiters.append({
            "name": result.get("title", "").split(" | ")[0],
            "title": result.get("title", ""),
            "link": result.get("link"),
            "company": company
        })

    return recruiters

def summarize_recruiters(recruiters):
    """Ranks recruiters using AI."""
    prompt = PromptTemplate(
        input_variables=["recruiters", "company"],
        template="""
        You must return JSON only. Do not include explanations.

        Rank the best recruiters for DataScience/ML/AI/Software roles at {company}.

        Output JSON:
        ```json
        [
            {{"rank": 1, "name": "Recruiter Name", "title": "Title", "link": "Profile URL", "insight": "Why they are ranked"}},
            {{"rank": 2, "name": "Recruiter Name", "title": "Title", "link": "Profile URL", "insight": "Why they are ranked"}},
            ...
        ]
        ```
        
        Recruiters:
        {recruiters}
        """
    )
    response = llm.invoke(prompt.format(recruiters=recruiters, company=recruiters[0]["company"])).content
    try:
        return json.loads(re.search(r"\[.*\]", response, re.DOTALL).group(0))
    except:
        return []

def generate_interview_questions(job_description):
    """Generates AI interview questions."""
    prompt = f"""
    Generate a **detailed** list of **technical and behavioral** interview questions for the following job description:
    
    {job_description}
    
    **STRICTLY return only JSON**. Do NOT include any extra text, headers, or explanations.
    
    Output Format:
    ```json
    {{
        "technical_questions": [
            "Question 1",
            "Question 2",
            ...
        ],
        "behavioral_questions": [
            "Question 1",
            "Question 2",
            ...
        ]
    }}
    ```
    """
    
    response = llm.invoke(prompt).content
    print(response)
    try:
        return json.loads(re.search(r"\{.*\}", response, re.DOTALL).group(0))
    except:
        return {"technical_questions": [], "behavioral_questions": []}

def evaluate_answer(question, user_answer):
    """Evaluates candidate response."""
    prompt = f"""
    Evaluate the response to: "{question}"

    Candidate's answer:
    {user_answer}

    JSON only:
    ```json
    {{
        "score": X,
        "feedback": "Feedback here"
    }}
    ```
    """
    
    response = llm.invoke(prompt).content
    try:
        return json.loads(re.search(r"\{.*\}", response, re.DOTALL).group(0))
    except:
        return {"score": "N/A", "feedback": "Could not evaluate response."}

# ===================== Streamlit UI ===================== #

st.set_page_config(page_title="Resume & Interview AI", layout="wide")
st.title("🚀 Resume Optimization & AI Interview Coach")

# Resume Upload Section
uploaded_resume = st.file_uploader("📂 Upload Resume (Word/PDF)", type=["docx", "pdf"])
job_url = st.text_input("🔗 Job Listing URL")

col1, col2 = st.columns(2)

with col1:
    if st.button("Fetch Job Details") and job_url:
        with st.spinner("Extracting job details..."):
            job_description, company_name = fetch_job_description(job_url)
        if job_description:
            st.session_state.job_description = job_description
            st.session_state.company_name = company_name
            st.success("✅ Job details fetched!")
    if "job_description" in st.session_state and st.session_state.job_description:
        st.text_area("📄 Job Description", st.session_state.job_description, height=250)

    if uploaded_resume and "job_description" in st.session_state:
        if st.button("Analyze Resume"):
            file_type = uploaded_resume.name.split(".")[-1]
            resume_text = extract_text_from_resume(uploaded_resume, file_type)
            
            if not resume_text:
                st.error("❌ Unsupported file type.")
            else:
                st.session_state.resume_text = resume_text
                st.subheader("🔎 Resume Analysis")
                st.write("✔ Your resume is analyzed. Missing skills will be highlighted here.")

with col2:
    if "company_name" in st.session_state:
        if st.button("Find Recruiters"):
            recruiters = summarize_recruiters(search_recruiters(st.session_state["company_name"]))
            if recruiters:
                for r in recruiters:
                    st.markdown(f"### [{r['name']}]({r['link']})")
                    st.write(r["title"])
                    st.markdown(f"[🔗 LinkedIn Profile]({r['link']})")
            else:
                st.warning("⚠ No recruiters found.")


if "job_description" in st.session_state and st.session_state.job_description:
    # Interview Chatbot Section
    st.subheader("🎤 AI-Powered Interview Practice")

    if "interview_questions" not in st.session_state:
        st.session_state.interview_questions = None

    if st.button("Generate Interview Questions") and "job_description" in st.session_state:
        with st.spinner("Generating questions..."):
            st.session_state.interview_questions = generate_interview_questions(st.session_state["job_description"])

    if "interview_questions" in st.session_state and st.session_state.interview_questions:
        st.subheader("📌 Technical Questions")
        for q in st.session_state.interview_questions["technical_questions"]:
            st.write(f"✔ {q}")

        st.subheader("💼 Behavioral Questions")
        for q in st.session_state.interview_questions["behavioral_questions"]:
            st.write(f"✔ {q}")

        st.subheader("📝 Answer & Get Feedback")
        question_selection = st.selectbox("Select a question", st.session_state.interview_questions["technical_questions"] + st.session_state.interview_questions["behavioral_questions"])
        # **New: Live Audio Recording**
        st.subheader("🎙 Record Your Answer")

        # Text input for the user's response
        user_answer = st.text_area("Your Answer", height=150)

        if st.button("Evaluate Answer"):
            with st.spinner("Evaluating your response..."):
                feedback = evaluate_answer(question_selection, user_answer)
            st.subheader("📊 Feedback")
            st.write(f"**Score:** {feedback['score']}/10")
            st.write(f"**Feedback:** {feedback['feedback']}")
